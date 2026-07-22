import json
from io import BytesIO

import matplotlib.pyplot as plt
import numpy as np
import pandas as pd
import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4, landscape
from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
from reportlab.lib.units import inch
from reportlab.platypus import (
    Image,
    PageBreak,
    Paragraph,
    SimpleDocTemplate,
    Spacer,
    Table,
    TableStyle,
)

from validation import (
    commitment_data_validation,
    spend_data_validation,
)
from reporting import (
    compare_statistical_methods,
    generate_cloud_report,
    suggest_method,
)


st.set_page_config(
    page_title="Overcast",
    page_icon="☁️",
    layout="wide",
)


#### Data & serialization helpers ####################################################
#### -> to turn in-memory program objects into JSON text format ######################

def make_json_safe(obj):
    """Converts NumPy and pandas values into JSON-safe Python values."""

    if isinstance(obj, dict):
        return {
            key: make_json_safe(value)
            for key, value in obj.items()
        }

    if isinstance(obj, list):
        return [
            make_json_safe(value)
            for value in obj
        ]

    if isinstance(obj, tuple):
        return [
            make_json_safe(value)
            for value in obj
        ]

    if isinstance(obj, pd.Timestamp):
        return obj.strftime("%Y-%m")

    if isinstance(obj, np.integer):
        return int(obj)

    if isinstance(obj, np.floating):
        return float(obj)

    if isinstance(obj, np.bool_):
        return bool(obj)

    return obj


def add_report_display_data(
    report,
    period_spend,
    period_start,
):
    """Adds historical and forecast dates needed in displaying / exporting a generated report."""

    report = report.copy()

    period_spend = period_spend.sort_values(
        "year_month"
    ).copy()

    has_history = len(period_spend) > 0

    if has_history:
        historical_dates = (
            period_spend["year_month"]
            .dt.strftime("%Y-%m")
            .tolist()
        )

        historical_spend = (
            period_spend["spend"]
            .astype(float)
            .tolist()
        )

        forecast_start = (
            period_spend["year_month"].max()
            + pd.DateOffset(months=1)
        )

    else:
        historical_dates = []
        historical_spend = []
        forecast_start = period_start

    forecast_dates = pd.date_range(
        start=forecast_start,
        periods=len(report["monthly_forecasts"]),
        freq="MS",
    ).strftime("%Y-%m").tolist()

    report["historical_dates"] = historical_dates
    report["historical_spend"] = historical_spend
    report["forecast_dates"] = forecast_dates

    return make_json_safe(report)


#### Plotting Helpers ###########################################################

def create_forecast_figure(report):
    """Create the historical and forecast monthly spend chart."""

    fig, ax = plt.subplots(figsize=(10, 5))

    if report["historical_dates"]:
        ax.plot(
            pd.to_datetime(report["historical_dates"]),
            report["historical_spend"],
            label="Historical",
        )

    if report["forecast_dates"]:
        ax.plot(
            pd.to_datetime(report["forecast_dates"]),
            report["monthly_forecasts"],
            label="Forecast",
        )

    ax.set_title(
        f"{report['cloud']} Historical and Forecasted Cloud Spend"
    )
    ax.set_xlabel("Month")
    ax.set_ylabel("Cost")
    ax.tick_params(axis="x", rotation=45)
    ax.grid(alpha=0.25)
    ax.legend()

    fig.tight_layout()

    return fig


def create_cumulative_figure(report):
    """Create the cumulative spend against commitment chart."""

    fig, ax = plt.subplots(figsize=(10, 5))

    historical_cumulative = (
        pd.Series(
            report["historical_spend"],
            dtype=float,
        )
        .cumsum()
        .tolist()
    )

    if report["historical_dates"]:
        ax.plot(
            pd.to_datetime(report["historical_dates"]),
            historical_cumulative,
            linewidth=2,
            label="Historical cumulative",
        )

        cumulative_start = historical_cumulative[-1]

    else:
        cumulative_start = 0.0

    forecast_cumulative = (
        cumulative_start
        + pd.Series(
            report["monthly_forecasts"],
            dtype=float,
        ).cumsum()
    ).tolist()

    if report["forecast_dates"]:
        ax.plot(
            pd.to_datetime(report["forecast_dates"]),
            forecast_cumulative,
            linewidth=2,
            label="Forecast cumulative",
        )

    ax.axhline(
        y=float(report["commitment"]),
        linestyle="--",
        linewidth=2,
        label="Commitment",
    )

    ax.set_title(
        f"{report['cloud']}: Cumulative Spend vs Commitment"
    )
    ax.set_xlabel("Month")
    ax.set_ylabel("Cumulative Spend")
    ax.tick_params(axis="x", rotation=45)
    ax.grid(alpha=0.25)
    ax.legend()

    fig.tight_layout()

    return fig


def figure_to_png(fig):
    """Convert a Matplotlib figure into an in-memory PNG."""

    image_buffer = BytesIO()

    fig.savefig(
        image_buffer,
        format="png",
        dpi=180,
        bbox_inches="tight",
    )

    image_buffer.seek(0)

    return image_buffer


#### PDF Export of Report #######################################################

def generate_report_pdf(report):
    """Generate a PDF version of the report."""

    pdf_buffer = BytesIO()

    document = SimpleDocTemplate(
        pdf_buffer,
        pagesize=landscape(A4),
        rightMargin=36,
        leftMargin=36,
        topMargin=30,
        bottomMargin=30,
    )

    styles = getSampleStyleSheet()

    title_style = ParagraphStyle(
        "ReportTitle",
        parent=styles["Title"],
        fontSize=22,
        leading=26,
        alignment=TA_CENTER,
        spaceAfter=16,
    )

    heading_style = ParagraphStyle(
        "ReportHeading",
        parent=styles["Heading2"],
        fontSize=14,
        leading=18,
        spaceBefore=10,
        spaceAfter=8,
    )

    story = []

    story.append(
        Paragraph(
            (
                f"Overcast: {report['cloud']} "
                "Cloud Cost Forecast"
            ),
            title_style,
        )
    )

    story.append(
        Paragraph(
            (
                f"<b>Commitment period:</b> "
                f"{report['period_start']} to "
                f"{report['period_end']}<br/>"
                f"<b>Forecasting method:</b> "
                f"{report['method_used']}<br/>"
                f"<b>Status:</b> "
                f"{report['status']}"
            ),
            styles["BodyText"],
        )
    )

    story.append(Spacer(1, 14))

    required_growth = report[
        "required_monthly_growth_rate (%)"
    ]

    required_growth_text = (
        "N/A"
        if required_growth is None
        else f"{float(required_growth):.2f}%"
    )

    metric_label_style = ParagraphStyle(
        "MetricLabel",
        parent=styles["BodyText"],
        fontSize=9,
        leading=11,
        textColor=colors.HexColor("#4A5568"),
        spaceAfter=4,
    )

    metric_value_style = ParagraphStyle(
        "MetricValue",
        parent=styles["BodyText"],
        fontSize=20,
        leading=23,
        textColor=colors.HexColor("#1F2937"),
    )

    detail_label_style = ParagraphStyle(
        "DetailLabel",
        parent=styles["BodyText"],
        fontSize=10,
        leading=13,
        textColor=colors.HexColor("#1F2937"),
    )

    required_growth = report["required_monthly_growth_rate (%)"]

    required_growth_text = (
        "N/A"
        if required_growth is None
        else f"{float(required_growth):.2f}%"
    )

    metric_cells = [
        [
            Paragraph(
                "Actual Cost to Date",
                metric_label_style,
            ),
            Paragraph(
                "Forecasted Total Spend",
                metric_label_style,
            ),
            Paragraph(
                "Commitment",
                metric_label_style,
            ),
        ],
        [
            Paragraph(
                f"${float(report['actual_cost_to_date']):,.2f}",
                metric_value_style,
            ),
            Paragraph(
                f"${float(report['forecasted_total_spend']):,.2f}",
                metric_value_style,
            ),
            Paragraph(
                f"${float(report['commitment']):,.2f}",
                metric_value_style,
            ),
        ],
        [
            Paragraph(
                "Gap",
                metric_label_style,
            ),
            Paragraph(
                "Future Spend Forecast",
                metric_label_style,
            ),
            Paragraph(
                "Required Monthly Growth",
                metric_label_style,
            ),
        ],
        [
            Paragraph(
                f"${float(report['gap']):,.2f}",
                metric_value_style,
            ),
            Paragraph(
                f"${float(report['future_spend_total']):,.2f}",
                metric_value_style,
            ),
            Paragraph(
                required_growth_text,
                metric_value_style,
            ),
        ],
    ]

    metric_layout = Table(
        metric_cells,
        colWidths=[
            3.25 * inch,
            3.25 * inch,
            3.25 * inch,
        ],
        rowHeights=[
            0.25 * inch,
            0.55 * inch,
            0.25 * inch,
            0.55 * inch,
        ],
    )

    metric_layout.setStyle(
        TableStyle(
            [
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("LEFTPADDING", (0, 0), (-1, -1), 0),
                ("RIGHTPADDING", (0, 0), (-1, -1), 18),
                ("TOPPADDING", (0, 0), (-1, -1), 0),
                ("BOTTOMPADDING", (0, 0), (-1, -1), 0),
            ]
        )
    )

    story.append(metric_layout)
    story.append(Spacer(1, 20))

    method_display = (
        str(report["method_used"])
        .replace("_forecast", "")
        .replace("_", " ")
        .title()
    )

    story.append(
        Paragraph(
            f"<b>Status:</b> {report['status']}",
            detail_label_style,
        )
    )

    story.append(Spacer(1, 10))

    story.append(
        Paragraph(
            f"<b>Method used:</b> {method_display}",
            detail_label_style,
        )
    )

    story.append(Spacer(1, 10))

    story.append(
        Paragraph(
            (
                f"<b>Period:</b> {report['period_start']} "
                f"to {report['period_end']}"
            ),
            detail_label_style,
        )
    )

    story.append(Spacer(1, 10))

    story.append(
        Paragraph(
            f"<b>Months remaining:</b> {report['months_remaining']}",
            detail_label_style,
        )
    )

    story.append(PageBreak())

    forecast_fig = create_forecast_figure(report)
    forecast_image = figure_to_png(forecast_fig)
    plt.close(forecast_fig)

    story.append(
        Paragraph(
            "Monthly Spend Forecast",
            heading_style,
        )
    )

    story.append(
        Image(
            forecast_image,
            width=9.7 * inch,
            height=4.8 * inch,
        )
    )

    story.append(PageBreak())

    cumulative_fig = create_cumulative_figure(report)
    cumulative_image = figure_to_png(cumulative_fig)
    plt.close(cumulative_fig)

    story.append(
        Paragraph(
            "Cumulative Spend vs Commitment",
            heading_style,
        )
    )

    story.append(
        Image(
            cumulative_image,
            width=9.7 * inch,
            height=4.8 * inch,
        )
    )

    story.append(PageBreak())

    story.append(
        Paragraph(
            "Monthly Forecast Values",
            heading_style,
        )
    )

    forecast_table_data = [
        [
            "Month",
            "Forecasted Spend",
        ]
    ]

    for month, forecast in zip(
        report["forecast_dates"],
        report["monthly_forecasts"],
    ):
        forecast_table_data.append(
            [
                month,
                f"${float(forecast):,.2f}",
            ]
        )

    forecast_table = Table(
        forecast_table_data,
        colWidths=[
            2.5 * inch,
            2.5 * inch,
        ],
        repeatRows=1,
    )

    forecast_table.setStyle(
        TableStyle(
            [
                (
                    "BACKGROUND",
                    (0, 0),
                    (-1, 0),
                    colors.HexColor("#DCE6F1"),
                ),
                (
                    "FONTNAME",
                    (0, 0),
                    (-1, 0),
                    "Helvetica-Bold",
                ),
                (
                    "ALIGN",
                    (1, 1),
                    (1, -1),
                    "RIGHT",
                ),
                (
                    "GRID",
                    (0, 0),
                    (-1, -1),
                    0.5,
                    colors.grey,
                ),
                (
                    "ROWBACKGROUNDS",
                    (0, 1),
                    (-1, -1),
                    [
                        colors.white,
                        colors.HexColor("#F5F5F5"),
                    ],
                ),
                (
                    "TOPPADDING",
                    (0, 0),
                    (-1, -1),
                    6,
                ),
                (
                    "BOTTOMPADDING",
                    (0, 0),
                    (-1, -1),
                    6,
                ),
            ]
        )
    )

    story.append(forecast_table)

    document.build(story)

    pdf_buffer.seek(0)

    return pdf_buffer.getvalue()


###### PowerPoint Export of Report ###################################################

def add_slide_title(slide, title):
    """Adds a title to a blank PowerPoint slide."""

    title_box = slide.shapes.add_textbox(
        Inches(0.6),
        Inches(0.25),
        Inches(12.1),
        Inches(0.7),
    )

    paragraph = title_box.text_frame.paragraphs[0]
    paragraph.text = title
    paragraph.font.size = Pt(26)
    paragraph.font.bold = True


def format_ppt_table_cell(
    cell,
    font_size=14,
    bold=False,
):
    """Formats a PowerPoint table cell."""

    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold


def add_monthly_forecast_slides(
    presentation,
    report,
    rows_per_slide=12,
):
    """Adds monthly forecast table slides in sections."""

    forecast_rows = list(
        zip(
            report["forecast_dates"],
            report["monthly_forecasts"],
        )
    )

    if not forecast_rows:
        return

    for start_index in range(
        0,
        len(forecast_rows),
        rows_per_slide,
    ):
        row_chunk = forecast_rows[
            start_index:start_index + rows_per_slide
        ]

        slide = presentation.slides.add_slide(
            presentation.slide_layouts[6]
        )

        if len(forecast_rows) > rows_per_slide:
            slide_number = (
                start_index // rows_per_slide
            ) + 1

            title = (
                "Monthly Forecast "
                f"(Part {slide_number})"
            )
        else:
            title = "Monthly Forecast"

        add_slide_title(slide, title)

        table_shape = slide.shapes.add_table(
            len(row_chunk) + 1,
            2,
            Inches(2.2),
            Inches(1.2),
            Inches(8.9),
            Inches(5.7),
        )

        table = table_shape.table

        table.columns[0].width = Inches(4.0)
        table.columns[1].width = Inches(4.9)

        table.cell(0, 0).text = "Month"
        table.cell(0, 1).text = "Forecasted Spend"

        format_ppt_table_cell(
            table.cell(0, 0),
            font_size=15,
            bold=True,
        )

        format_ppt_table_cell(
            table.cell(0, 1),
            font_size=15,
            bold=True,
        )

        for row_index, (
            month,
            forecast,
        ) in enumerate(
            row_chunk,
            start=1,
        ):
            table.cell(
                row_index,
                0,
            ).text = str(month)

            table.cell(
                row_index,
                1,
            ).text = f"${float(forecast):,.2f}"

            format_ppt_table_cell(
                table.cell(row_index, 0),
                font_size=13,
            )

            format_ppt_table_cell(
                table.cell(row_index, 1),
                font_size=13,
            )


def generate_report_pptx(report):
    """Generates a PowerPoint version of the report."""

    presentation = Presentation()

    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    # Slide 1: Title
    title_slide = presentation.slides.add_slide(presentation.slide_layouts[0])

    title = title_slide.shapes.title
    title.text = (
        f"Overcast: {report['cloud']} Cloud Cost Forecast"
    )

    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    subtitle = title_slide.placeholders[1]
    subtitle.text = (
        f"{report['period_start']} to {report['period_end']}\n"
        f"Method: {method_display}"
    )

    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER

    title.text_frame.paragraphs[0].font.size = Pt(30)

    for paragraph in subtitle.text_frame.paragraphs:
        paragraph.font.size = Pt(18)

    # Slide 2: Forecast Summary 
    
    summary_slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_slide_title(
        summary_slide,
        "Forecast Summary",
    )

    required_growth = report[
        "required_monthly_growth_rate (%)"
    ]

    required_growth_text = (
        "N/A"
        if required_growth is None
        else f"{float(required_growth):.2f}%"
    )

    method_display = (
        str(report["method_used"])
        .replace("_forecast", "")
        .replace("_", " ")
        .title()
    )

    summary_metrics = [
        (
            "Actual Cost to Date",
            f"${float(report['actual_cost_to_date']):,.2f}",
        ),
        (
            "Forecasted Total Spend",
            f"${float(report['forecasted_total_spend']):,.2f}",
        ),
        (
            "Commitment",
            f"${float(report['commitment']):,.2f}",
        ),
        (
            "Gap",
            f"${float(report['gap']):,.2f}",
        ),
        (
            "Future Spend Forecast",
            f"${float(report['future_spend_total']):,.2f}",
        ),
        (
            "Required Monthly Growth",
            required_growth_text,
        ),
    ]

    metric_positions = [
        (0.6, 1.15),
        (4.75, 1.15),
        (8.9, 1.15),
        (0.6, 2.75),
        (4.75, 2.75),
        (8.9, 2.75),
    ]

    for (
        label,
        value,
    ), (
        left,
        top,
    ) in zip(
        summary_metrics,
        metric_positions,
    ):
        label_box = summary_slide.shapes.add_textbox(
            Inches(left),
            Inches(top),
            Inches(3.7),
            Inches(0.35),
        )

        label_paragraph = (
            label_box.text_frame.paragraphs[0]
        )

        label_paragraph.text = label
        label_paragraph.font.size = Pt(11)

        value_box = summary_slide.shapes.add_textbox(
            Inches(left),
            Inches(top + 0.35),
            Inches(3.7),
            Inches(0.65),
        )

        value_paragraph = (
            value_box.text_frame.paragraphs[0]
        )

        value_paragraph.text = value
        value_paragraph.font.size = Pt(23)
        value_paragraph.font.bold = False

    details_box = summary_slide.shapes.add_textbox(
        Inches(0.6),
        Inches(4.35),
        Inches(12.0),
        Inches(2.4),
    )

    details_frame = details_box.text_frame
    details_frame.clear()

    status_paragraph = details_frame.paragraphs[0]
    status_paragraph.text = (
        f"Status: {report['status']}"
    )
    status_paragraph.font.size = Pt(13)

    method_paragraph = details_frame.add_paragraph()
    method_paragraph.text = (
        f"Method used: {method_display}"
    )
    method_paragraph.font.size = Pt(13)
    method_paragraph.space_before = Pt(10)

    period_paragraph = details_frame.add_paragraph()
    period_paragraph.text = (
        f"Period: {report['period_start']} "
        f"to {report['period_end']}"
    )
    period_paragraph.font.size = Pt(13)
    period_paragraph.space_before = Pt(10)

    months_paragraph = details_frame.add_paragraph()
    months_paragraph.text = (
        f"Months remaining: {report['months_remaining']}"
    )
    months_paragraph.font.size = Pt(13)
    months_paragraph.space_before = Pt(10)

    # Slide 3: Monthly Forecast Chart

    forecast_slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_slide_title(
        forecast_slide,
        "Historical and Forecasted Spend",
    )

    forecast_fig = create_forecast_figure(report)
    forecast_image = figure_to_png(forecast_fig)
    plt.close(forecast_fig)

    forecast_slide.shapes.add_picture(
        forecast_image,
        Inches(0.9),
        Inches(1.1),
        width=Inches(11.6),
        height=Inches(5.8),
    )

    # Slide 4: Cumulative Forecast Chart
    cumulative_slide = presentation.slides.add_slide(
        presentation.slide_layouts[6]
    )

    add_slide_title(
        cumulative_slide,
        "Cumulative Spend vs Commitment",
    )

    cumulative_fig = create_cumulative_figure(
        report
    )

    cumulative_image = figure_to_png(
        cumulative_fig
    )

    plt.close(cumulative_fig)

    cumulative_slide.shapes.add_picture(
        cumulative_image,
        Inches(0.9),
        Inches(1.1),
        width=Inches(11.6),
        height=Inches(5.8),
    )

    # Slides 5-End: Monthly Forecast Tables
    add_monthly_forecast_slides(
        presentation=presentation,
        report=report,
        rows_per_slide=12,
    )

    pptx_buffer = BytesIO()

    presentation.save(pptx_buffer)

    pptx_buffer.seek(0)

    return pptx_buffer.getvalue()


#### Streamlit Report Display ##################################


def display_report(report):
    """Display a generated or loaded report in Streamlit."""

    col1, col2, col3 = st.columns(3)

    col1.metric(
        "Actual Cost to Date",
        f"${float(report['actual_cost_to_date']):,.2f}",
    )

    col2.metric(
        "Forecasted Total Spend",
        f"${float(report['forecasted_total_spend']):,.2f}",
    )

    col3.metric(
        "Commitment",
        f"${float(report['commitment']):,.2f}",
    )

    col4, col5, col6 = st.columns(3)

    col4.metric(
        "Gap",
        f"${float(report['gap']):,.2f}",
    )

    col5.metric(
        "Future Spend Forecast",
        f"${float(report['future_spend_total']):,.2f}",
    )

    required_growth = report[
        "required_monthly_growth_rate (%)"
    ]

    col6.metric(
        "Required Monthly Growth",
        (
            "N/A"
            if required_growth is None
            else f"{float(required_growth):.2f}%"
        ),
    )

    st.write(
        f"**Status:** {report['status']}"
    )

    st.write(
        f"**Method used:** {report['method_used']}"
    )

    st.write(
        (
            f"**Period:** {report['period_start']} "
            f"to {report['period_end']}"
        )
    )

    st.write(
        f"**Months remaining:** "
        f"{report['months_remaining']}"
    )

    st.subheader("Monthly Forecast")

    forecast_df = pd.DataFrame(
        {
            "Month": report["forecast_dates"],
            "Forecasted Spend": report[
                "monthly_forecasts"
            ],
        }
    )

    st.dataframe(
        forecast_df,
        use_container_width=True,
        hide_index=True,
    )

    st.subheader("Forecast Plot")

    forecast_fig = create_forecast_figure(
        report
    )

    st.pyplot(forecast_fig)

    plt.close(forecast_fig)

    st.subheader(
        "Cumulative Spend vs Commitment"
    )

    cumulative_fig = create_cumulative_figure(
        report
    )

    st.pyplot(cumulative_fig)

    plt.close(cumulative_fig)

    with st.expander("View raw report data"):
        st.json(report)


def display_download_buttons(
    report,
    cloud_name,
):
    """Display JSON, PDF, and PowerPoint download buttons."""

    report_json = json.dumps(
        report,
        indent=4,
    )

    try:
        report_pdf = generate_report_pdf(
            report
        )
    except Exception as error:
        report_pdf = None

        st.error(
            f"PDF generation failed: {error}"
        )

    try:
        report_pptx = generate_report_pptx(
            report
        )
    except Exception as error:
        report_pptx = None

        st.error(
            f"PowerPoint generation failed: {error}"
        )

    st.subheader("Download Report")

    download_col1, download_col2, download_col3 = (
        st.columns(3)
    )

    download_col1.download_button(
        label="Download JSON",
        data=report_json,
        file_name=(
            f"Overcast_{cloud_name}_report.json"
        ),
        mime="application/json",
        use_container_width=True,
    )

    if report_pdf is not None:
        download_col2.download_button(
            label="Download PDF",
            data=report_pdf,
            file_name=(
                f"Overcast_{cloud_name}_report.pdf"
            ),
            mime="application/pdf",
            use_container_width=True,
        )

    if report_pptx is not None:
        download_col3.download_button(
            label="Download PowerPoint",
            data=report_pptx,
            file_name=(
                f"Overcast_{cloud_name}_report.pptx"
            ),
            mime=(
                "application/vnd.openxmlformats-"
                "officedocument.presentationml."
                "presentation"
            ),
            use_container_width=True,
        )


### Session state ########################################################

if "generated_report" not in st.session_state:
    st.session_state.generated_report = None

if "generated_report_cloud" not in st.session_state:
    st.session_state.generated_report_cloud = None


### Application interface (UI) ############################################

st.title("Overcast")
st.caption("A cloud cost forecasting tool.")


# Load report section
st.subheader("Load Saved Report")

saved_report_file = st.file_uploader(
    "Upload saved report JSON",
    type=["json"],
    key="saved_report_uploader",
)

if saved_report_file is not None:
    try:
        saved_report = json.load(
            saved_report_file
        )

        st.success(
            "Saved report loaded successfully."
        )

        display_report(saved_report)

        loaded_cloud = saved_report.get(
            "cloud",
            "Cloud",
        )

        display_download_buttons(
            report=saved_report,
            cloud_name=loaded_cloud,
        )

    except (
        json.JSONDecodeError,
        KeyError,
        TypeError,
        ValueError,
    ) as error:
        st.error(
            f"Unable to load the report: {error}"
        )


# Upload data section
st.subheader("1. Upload Data")

spend_file = st.file_uploader(
    "Upload spend data CSV",
    type=["csv"],
    key="spend_uploader",
)

commitment_file = st.file_uploader(
    "Upload commitment data CSV",
    type=["csv"],
    key="commitment_uploader",
)


if (
    spend_file is not None
    and commitment_file is not None
):
    try:
        spend_data = spend_data_validation(
            pd.read_csv(spend_file)
        )

        commitments_data = (
            commitment_data_validation(
                pd.read_csv(commitment_file)
            )
        )

    except Exception as error:
        st.error(
            f"Data validation failed: {error}"
        )

        st.stop()

    st.success("Data loaded successfully.")

    with st.expander(
        "View uploaded spend data"
    ):
        st.dataframe(
            spend_data,
            use_container_width=True,
        )

    with st.expander(
        "View uploaded commitment data"
    ):
        st.dataframe(
            commitments_data,
            use_container_width=True,
        )

    selected_cloud = st.selectbox(
        "Select cloud provider",
        commitments_data["cloud"].unique(),
    )

    cloud_commitment = commitments_data[
        commitments_data["cloud"]
        == selected_cloud
    ].iloc[0]

    period_start = cloud_commitment[
        "period_start"
    ]

    period_end = cloud_commitment[
        "period_end"
    ]

    period_spend = spend_data[
        (
            spend_data["cloud"]
            == selected_cloud
        )
        & (
            spend_data["year_month"]
            >= period_start
        )
        & (
            spend_data["year_month"]
            <= period_end
        )
    ].sort_values("year_month")

    has_historical_data = (
        len(period_spend) > 0
    )

    selected_method = None
    project_adjustments = []

    st.subheader("2. Forecast Setup")

    if has_historical_data:
        st.write(
            (
                "Historical spend data was found "
                "for this cloud and commitment period."
            )
        )

        st.subheader(
            "Forecast Method Comparison"
        )

        try:
            comparison_df = (
                compare_statistical_methods(
                    spend_data=spend_data,
                    commitments=commitments_data,
                    cloud=selected_cloud,
                )
            )

            st.dataframe(
                comparison_df,
                use_container_width=True,
                hide_index=True,
            )

        except Exception as error:
            st.warning(
                (
                    "The forecast comparison could "
                    f"not be generated: {error}"
                )
            )

        spend_series = period_spend[
            "spend"
        ]

        has_seasonality = st.checkbox(
            (
                "Use seasonal forecasting "
                "if data appears seasonal"
            ),
            value=False,
            help=(
                "Allows Holt-Winters forecasting "
                "when at least 24 months of data "
                "are available. Use this when a "
                "seasonal pattern is expected."
            ),
        )

        suggestion = suggest_method(
            spend_series,
            has_seasonality=has_seasonality,
        )

        st.subheader("Suggested Method")

        st.write(
            (
                "**Suggested method:** "
                f"{suggestion['suggested_method']}"
            )
        )

        st.write(
            f"**Reason:** {suggestion['reason']}"
        )

        with st.expander("View diagnostics"):
            st.json(
                suggestion["diagnostics"]
            )

        method_options = [
            "run_rate_forecast",
            "moving_average_forecast",
            "historic_growth_forecast",
            (
                "single_exponential_"
                "smoothing_forecast"
            ),
            (
                "holt_double_exponential_"
                "smoothing_forecast"
            ),
            (
                "holtwinters_triple_"
                "exponential_smoothing_forecast"
            ),
        ]

        default_index = (
            method_options.index(
                suggestion["suggested_method"]
            )
            if suggestion["suggested_method"]
            in method_options
            else 0
        )

        selected_method = st.selectbox(
            "Choose method for final report",
            method_options,
            index=default_index,
        )

        apply_project_adjustments = (
            st.checkbox(
                "Apply planned project adjustment"
            )
        )

        if apply_project_adjustments:
            latest_spend_month = (
                period_spend["year_month"].max()
            )

            months_remaining_preview = (
                (
                    period_end.year
                    - latest_spend_month.year
                )
                * 12
                + (
                    period_end.month
                    - latest_spend_month.month
                )
            )

            if months_remaining_preview > 0:
                future_months = pd.date_range(
                    start=(
                        latest_spend_month
                        + pd.DateOffset(months=1)
                    ),
                    periods=(
                        months_remaining_preview
                    ),
                    freq="MS",
                )

                adjustment_df = pd.DataFrame(
                    {
                        "month": (
                            future_months.strftime(
                                "%Y-%m"
                            )
                        ),
                        "project_adjustment": (
                            [0.0]
                            * months_remaining_preview
                        ),
                    }
                )

                edited_adjustment_df = (
                    st.data_editor(
                        adjustment_df,
                        use_container_width=True,
                        disabled=["month"],
                        hide_index=True,
                    )
                )

                project_adjustments = (
                    pd.to_numeric(
                        edited_adjustment_df[
                            "project_adjustment"
                        ],
                        errors="coerce",
                    )
                    .fillna(0.0)
                    .tolist()
                )

            else:
                st.info(
                    (
                        "There are no future months "
                        "remaining in the commitment "
                        "period."
                    )
                )

    else:
        st.warning(
            (
                "No historical spend data was found "
                "for this cloud and commitment "
                "period. At least one month of "
                "historical spend data is required "
                "to generate a forecast."
            )
        )

    st.subheader("3. Final Report")

    generate_disabled = (
        not has_historical_data
        or selected_method is None
    )

    if st.button(
        "Generate Report",
        disabled=generate_disabled):
        try:
            report = generate_cloud_report(
                spend_data=spend_data,
                commitments=commitments_data,
                cloud=selected_cloud,
                selected_method=selected_method,
                project_adjustments=(
                    project_adjustments
                ),
            )

        except Exception as error:
            st.error(
                (
                    "The report could not be "
                    f"generated: {error}"
                )
            )

            report = None

        if report is not None:
            failed_statuses = [
                "Insufficient data",
                "Forecast failed",
            ]

            if (
                "status" in report
                and report["status"]
                in failed_statuses
            ):
                st.error(
                    report.get(
                        "message",
                        "The forecast failed.",
                    )
                )

                st.json(report)

            else:
                report = add_report_display_data(
                    report=report,
                    period_spend=period_spend,
                    period_start=period_start,
                )

                st.session_state[
                    "generated_report"
                ] = report

                st.session_state[
                    "generated_report_cloud"
                ] = selected_cloud

    if (
        st.session_state.generated_report
        is not None
    ):
        current_report = (
            st.session_state.generated_report
        )

        current_cloud = (
            st.session_state[
                "generated_report_cloud"
            ]
            or current_report.get(
                "cloud",
                "Cloud",
            )
        )

        display_report(current_report)

        display_download_buttons(
            report=current_report,
            cloud_name=current_cloud,
        )